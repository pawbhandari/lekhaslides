from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import logging
from typing import List

logger = logging.getLogger("lekhaslides.pptx")

def create_pptx_from_images(images: List[Image.Image]) -> io.BytesIO:
    """
    Create PPTX file from list of PIL images
    
    Args:
        images: List of PIL Image objects (1920x1080)
    
    Returns:
        BytesIO buffer containing PPTX file
    """
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for idx, img in enumerate(images):
        try:
            # Ensure image is in RGB mode (PPTX doesn't support RGBA or other modes well)
            if img.mode != 'RGB':
                print(f"Converting slide {idx+1} from {img.mode} to RGB")
                img = img.convert('RGB')
            
            # Convert PIL image to bytes
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)
            
            # Add blank slide
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            
            # Add image filling entire slide
            slide.shapes.add_picture(
                img_byte_arr,
                Inches(0),
                Inches(0),
                width=prs.slide_width,
                height=prs.slide_height
            )
            logger.info(f"Added slide {idx+1}/{len(images)}")
        except Exception as e:
            logger.error(f"Error adding slide {idx+1}: {str(e)}")
            raise
    
    # Save to BytesIO
    pptx_output = io.BytesIO()
    prs.save(pptx_output)
    pptx_output.seek(0)
    
    return pptx_output
